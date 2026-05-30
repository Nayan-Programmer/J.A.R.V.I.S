"""
PPT SERVICE
===========
Generates PowerPoint presentations using python-pptx.
LLM (Groq) creates the slide outline; python-pptx builds the file.
Install: pip install python-pptx lxml
"""

import json
import logging
import re
from pathlib import Path
from typing import Optional

logger = logging.getLogger("J.A.R.V.I.S")

GENERATED_DIR = Path(__file__).parent.parent.parent / "database" / "generated"
GENERATED_DIR.mkdir(parents=True, exist_ok=True)

THEME = {
    "bg":         "1a1a2e",
    "title_text": "e2e2f0",
    "body_text":  "c9c9dd",
    "accent":     "7c5cbf",
}


def _hex(h):
    from pptx.dml.color import RGBColor
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def generate_ppt(topic: str, groq_service) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    outline = _get_outline(topic, groq_service)
    slides_data = outline.get("slides", [])
    if not slides_data:
        raise ValueError("LLM returned no slides for the presentation.")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, sinfo in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, THEME["bg"])

        title  = (sinfo.get("title") or "").strip()
        points = sinfo.get("points") or []
        note   = (sinfo.get("note") or "").strip()

        if idx == 0:
            _textbox(slide, title,
                     Inches(1), Inches(2.5), Inches(11.33), Inches(1.4),
                     Pt(44), True, THEME["title_text"], PP_ALIGN.CENTER)
            if note:
                _textbox(slide, note,
                         Inches(2), Inches(4.2), Inches(9.33), Inches(0.8),
                         Pt(22), False, THEME["body_text"], PP_ALIGN.CENTER)
            _rect(slide, 0, Inches(4.05), prs.slide_width, Inches(0.07), THEME["accent"])
        else:
            _rect(slide, 0, 0, prs.slide_width, Inches(1.25), THEME["accent"])
            _textbox(slide, title,
                     Inches(0.35), Inches(0.15), Inches(12.63), Inches(1.0),
                     Pt(30), True, THEME["title_text"], PP_ALIGN.LEFT)
            top = Inches(1.5)
            for pt in points[:7]:
                _textbox(slide, "•  " + pt,
                         Inches(0.6), top, Inches(12.13), Inches(0.65),
                         Pt(19), False, THEME["body_text"], PP_ALIGN.LEFT)
                top += Inches(0.72)
            _textbox(slide, str(idx + 1),
                     Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.35),
                     Pt(12), False, THEME["body_text"], PP_ALIGN.RIGHT)

    safe = re.sub(r"[^\w\s-]", "", topic.lower()).strip().replace(" ", "_")[:40]
    out = GENERATED_DIR / (safe + ".pptx")
    prs.save(str(out))
    logger.info("[PPT] Saved: %s", out)
    return out


def _get_outline(topic: str, groq_service) -> dict:
    prompt = (
        "Create a professional PowerPoint outline for: \"" + topic + "\"\n\n"
        "Return ONLY valid JSON, no markdown, no code fences:\n"
        "{\n"
        '  "slides": [\n'
        '    {"title": "Cover Title", "note": "subtitle"},\n'
        '    {"title": "Slide Title", "points": ["Point 1", "Point 2", "Point 3"]},\n'
        '    {"title": "Conclusion", "points": ["Key takeaway 1", "Key takeaway 2"]}\n'
        "  ]\n"
        "}\n\n"
        "Rules:\n"
        "- First slide is cover: title + note (no points)\n"
        "- Other slides: title + 4-6 bullet points (under 12 words each)\n"
        "- 8-12 slides total. Return ONLY the JSON object."
    )
    raw = groq_service.get_response(prompt, chat_history=None)
    raw = raw.strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.M)
    raw = re.sub(r"```\s*$", "", raw, flags=re.M).strip()
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        m = re.search(r"\{[\s\S]+\}", raw)
        if m:
            return json.loads(m.group(0))
        raise ValueError("Could not parse LLM outline: " + raw[:200])


def _set_bg(slide, hex_colour: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex(hex_colour)


def _rect(slide, left, top, width, height, hex_colour: str):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(hex_colour)
    shape.line.fill.background()


def _textbox(slide, text, left, top, width, height, font_size, bold, colour, align):
    from pptx.enum.text import PP_ALIGN
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = font_size
    run.font.bold      = bold
    run.font.color.rgb = _hex(colour)
    run.font.name      = "Calibri"
